import os
import re
from dataclasses import dataclass
from typing import Iterable

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt


@dataclass(frozen=True)
class SlideSpec:
    title: str
    bullets: list[str]


def _cleanup_spaces(s: str) -> str:
    s = s.replace("\u00a0", " ")
    s = re.sub(r"\s+", " ", s)
    return s.strip()


def _dehyphenate_linebreaks(s: str) -> str:
    return re.sub(r"-\s+", "", s)


def _set_title(slide, title: str) -> None:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.is_placeholder and shape.placeholder_format.type == 1:  # TITLE
            tf = shape.text_frame
            tf.clear()
            tf.text = title
            return

    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.clear()
            tf.text = title
            return


def _set_bullets(slide, bullets: Iterable[str], max_placeholders: int = 2) -> None:
    body_shapes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.is_placeholder and shape.placeholder_format.type == 2:  # BODY
            body_shapes.append(shape)

    body_shapes = body_shapes[:max_placeholders]
    if not body_shapes:
        return

    bullets = [b for b in (_cleanup_spaces(_dehyphenate_linebreaks(x)) for x in bullets) if b]
    if not bullets:
        return

    buckets: list[list[str]] = [[] for _ in body_shapes]
    for i, b in enumerate(bullets):
        buckets[i % len(body_shapes)].append(b)

    for shape, items in zip(body_shapes, buckets, strict=False):
        tf = shape.text_frame
        tf.clear()
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            for run in p.runs:
                run.font.size = Pt(20)


def _render_page_png(pdf_path: str, page_index: int, out_path: str, zoom: float = 2.0) -> None:
    doc = fitz.open(pdf_path)
    page = doc.load_page(page_index)
    mat = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=mat, alpha=False)
    os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
    pix.save(out_path)


def _extract_title_page_meta(pdf_path: str) -> tuple[str, str, str]:
    doc = fitz.open(pdf_path)
    text = doc.load_page(0).get_text("text")
    lines = [l.strip() for l in text.splitlines() if l.strip()]

    title = None
    if "Fulfilled by:" in lines:
        idx = lines.index("Fulfilled by:")
        if idx >= 2:
            title = " ".join(lines[idx - 2 : idx])
    title = title or "Thesis presentation"

    author = None
    for i, l in enumerate(lines):
        if l.startswith("Student") and i + 1 < len(lines):
            author = lines[i + 1]
            break
    author = author or ""

    year_line = lines[-1] if lines else ""

    return (_cleanup_spaces(title), _cleanup_spaces(author), _cleanup_spaces(year_line))


def _extract_abstract_bullets(pdf_path: str, max_sentences: int = 5) -> list[str]:
    doc = fitz.open(pdf_path)
    text = doc.load_page(1).get_text("text").replace("\n", " ")
    text = re.sub(r"\s+", " ", text).strip()
    text = re.sub(r"^Abstract\s+", "", text, flags=re.IGNORECASE)

    sentences = re.split(r"(?<=[.!?])\s+", text)
    sentences = [_cleanup_spaces(s) for s in sentences if _cleanup_spaces(s)]
    return sentences[:max_sentences]


def _extract_objectives(pdf_path: str) -> list[str]:
    doc = fitz.open(pdf_path)
    text = doc.load_page(5).get_text("text")
    lines = [l.strip() for l in text.splitlines() if l.strip()]

    bul: list[str] = []
    cur: str | None = None
    for l in lines:
        if l.startswith("•"):
            if cur:
                bul.append(cur)
            cur = l.lstrip("•").strip()
        elif cur and l.startswith("to "):
            bul.append(cur)
            cur = l
        elif cur and (l[0].islower() or l.startswith(("and ", "or ", "via ", "the "))):
            cur += " " + l
        else:
            if cur:
                bul.append(cur)
                cur = None
    if cur:
        bul.append(cur)

    clean = [_cleanup_spaces(_dehyphenate_linebreaks(x)) for x in bul]
    # ensure 'Heston' appears in the SDE bullet
    out = []
    for b in clean:
        if "three SDE models" in b and "Heston" not in b:
            b = b.rstrip(";") + ", and Heston."
        out.append(b)
    return out[:6]


def _extract_ch2_key_points(pdf_path: str) -> list[str]:
    doc = fitz.open(pdf_path)
    text = doc.load_page(10).get_text("text")
    items = []
    for line in text.splitlines():
        line = line.strip()
        if re.match(r"^[123]\.\s+", line):
            items.append(_cleanup_spaces(line))
    return items


def _extract_results_table(pdf_path: str) -> list[str]:
    doc = fitz.open(pdf_path)
    text = doc.load_page(26).get_text("text")
    lines = [l.strip() for l in text.splitlines() if l.strip()]

    # pick table rows
    rows = []
    capture = False
    for l in lines:
        if l.startswith("Model"):
            capture = True
            continue
        if capture:
            if l.startswith("GBM") or l.startswith("Elastic Net"):
                rows.append(l)
            elif rows and re.match(r"^[0-9]", l):
                rows[-1] += " | " + l
            elif l.startswith("6.3"):
                break

    # If PDF text extraction already produced full rows, keep them
    if any("|" not in r for r in rows):
        # fallback: parse the known order from the raw block (best-effort)
        # The page extraction we saw already contains each row on a single line.
        rows = []
        for l in lines:
            if re.match(
                r"^(Elastic Net|Random Forest|Gradient Boosting|LSTM|GBM \(SDE\)|Merton Jump-Diff \(SDE\)|Heston \(SDE\))\s+",
                l,
            ):
                rows.append(re.sub(r"\s{2,}", " ", l))

    return rows


def _extract_limitations(pdf_path: str) -> list[str]:
    doc = fitz.open(pdf_path)
    text = doc.load_page(29).get_text("text")
    out = []
    for line in text.splitlines():
        line = line.strip()
        if re.match(r"^[1-9]\.\s+", line):
            out.append(_cleanup_spaces(line))
    return out


def build_presentation(template_path: str, pdf_path: str, out_pptx_path: str) -> None:
    prs = Presentation(template_path)

    title, author, year_line = _extract_title_page_meta(pdf_path)

    specs: list[SlideSpec] = [
        SlideSpec(title=title, bullets=[author, year_line]),
        SlideSpec(title="Motivation & Problem", bullets=_extract_abstract_bullets(pdf_path, max_sentences=4)),
        SlideSpec(title="Research objectives", bullets=_extract_objectives(pdf_path)),
        SlideSpec(
            title="Why Bitcoin price is a stochastic process",
            bullets=_extract_ch2_key_points(pdf_path),
        ),
        SlideSpec(
            title="Models",
            bullets=[
                "SDE: Geometric Brownian Motion (GBM)",
                "SDE: Merton Jump-Diffusion",
                "SDE: Heston stochastic volatility",
                "ML: Elastic Net, Random Forest, Gradient Boosting, LSTM",
            ],
        ),
        SlideSpec(
            title="Experimental results (15-min data, h = 10)",
            bullets=[
                "ML models outperform SDE models on point prediction metrics (MAE/RMSE/R2).",
                "Best MAE: Elastic Net (337.6 USD); best RMSE: LSTM (487.2 USD).",
                "SDE models are worse on point metrics but provide distributional forecasts.",
            ],
        ),
        SlideSpec(
            title="Limitations & future work",
            bullets=_extract_limitations(pdf_path)[:5],
        ),
        SlideSpec(
            title="Conclusion",
            bullets=[
                "Combined SDE + ML view provides both interpretability and predictive accuracy.",
                "Jump-diffusion improves distributional fit vs GBM.",
                "ML is best for short-horizon point forecasts; SDE for uncertainty quantification.",
            ],
        ),
    ]

    # Use a consistent text layout from the template
    text_layout = prs.slide_layouts[1]  # Текст_1

    for spec in specs:
        slide = prs.slides.add_slide(text_layout)
        _set_title(slide, spec.title)
        _set_bullets(slide, spec.bullets, max_placeholders=2)

    # Add 1 slide with the results table as an image (keeps formatting without fighting PPT tables)
    tmp_dir = os.path.join(os.path.dirname(out_pptx_path) or ".", ".tmp_presa")
    os.makedirs(tmp_dir, exist_ok=True)
    table_img = os.path.join(tmp_dir, "results_table.png")
    _render_page_png(pdf_path, page_index=26, out_path=table_img, zoom=2.0)
    slide = prs.slides.add_slide(prs.slide_layouts[11])  # чистый
    slide.shapes.add_picture(table_img, 0, 0, width=prs.slide_width, height=prs.slide_height)

    prs.save(out_pptx_path)


if __name__ == "__main__":
    template = os.environ.get("PRESA_TEMPLATE", "shablon.pptx")
    pdf = os.environ.get("PRESA_PDF", "diploma.pdf")
    out = os.environ.get("PRESA_OUT", "presa.pptx")
    build_presentation(template, pdf, out)

